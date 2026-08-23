"""Deterministic document extraction, PII masking and traceable semantic chunks."""

import csv
import hashlib
import io
import json
import re
import zipfile
from dataclasses import dataclass
from html.parser import HTMLParser
from pathlib import PurePath
from typing import Protocol
from urllib.parse import unquote, urlsplit
from xml.etree import ElementTree

from docx import Document as DocxDocument
from pypdf import PdfReader
from pptx import Presentation

from blogops.core.errors import AppError

ALLOWED_UPLOADS: dict[str, frozenset[str]] = {
    ".pdf": frozenset({"application/pdf"}),
    ".docx": frozenset({"application/vnd.openxmlformats-officedocument.wordprocessingml.document"}),
    ".pptx": frozenset({"application/vnd.openxmlformats-officedocument.presentationml.presentation"}),
    ".txt": frozenset({"text/plain"}),
    ".md": frozenset({"text/markdown", "text/plain"}),
    ".csv": frozenset({"text/csv", "application/csv", "text/plain"}),
}

_EMAIL = re.compile(r"(?<![\w.+-])[\w.+-]+@[\w-]+(?:\.[\w-]+)+")
_PHONE = re.compile(r"(?<!\d)(?:\+?82[- ]?)?0?1[016789][ -]?\d{3,4}[ -]?\d{4}(?!\d)")
_RESIDENT_ID = re.compile(r"(?<!\d)\d{6}[ -]?[1-4]\d{6}(?!\d)")


@dataclass(frozen=True, slots=True)
class ExtractedBlock:
    text: str
    locator: dict[str, int | str]
    block_type: str = "paragraph"


@dataclass(frozen=True, slots=True)
class ParsedDocument:
    title: str
    blocks: tuple[ExtractedBlock, ...]
    parser_name: str
    parser_version: str = "1"
    quality_score: float = 1.0

    @property
    def text(self) -> str:
        return "\n\n".join(block.text for block in self.blocks if block.text)


@dataclass(frozen=True, slots=True)
class ChunkDraft:
    sequence: int
    text: str
    locator: dict[str, int | str]
    text_hash: str
    token_estimate: int
    pii_masked: bool


class Parser(Protocol):
    def parse(self, filename: str, data: bytes) -> ParsedDocument: ...


def validate_upload_metadata(filename: str, content_type: str, size: int, max_size: int) -> str:
    suffix = PurePath(filename).suffix.lower()
    if suffix not in ALLOWED_UPLOADS:
        raise AppError("FILE_TYPE_NOT_ALLOWED", "지원하지 않는 파일 형식입니다.", 422)
    normalized_type = content_type.split(";", 1)[0].strip().lower()
    if normalized_type not in ALLOWED_UPLOADS[suffix]:
        raise AppError("FILE_MIME_MISMATCH", "파일 확장자와 MIME 형식이 일치하지 않습니다.", 422)
    if size <= 0 or size > max_size:
        raise AppError("FILE_SIZE_INVALID", "파일 크기가 허용 범위를 벗어났습니다.", 422)
    return suffix


def validate_file_signature(suffix: str, data: bytes) -> None:
    if suffix == ".pdf" and not data.startswith(b"%PDF-"):
        raise AppError("FILE_SIGNATURE_MISMATCH", "PDF 파일 서명이 올바르지 않습니다.", 422)
    if suffix in {".docx", ".pptx"}:
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as archive:
                names = frozenset(archive.namelist())
        except zipfile.BadZipFile as exc:
            raise AppError("FILE_SIGNATURE_MISMATCH", "Office 파일 구조가 올바르지 않습니다.", 422) from exc
        required_prefix = "word/" if suffix == ".docx" else "ppt/"
        if "[Content_Types].xml" not in names or not any(name.startswith(required_prefix) for name in names):
            raise AppError("FILE_SIGNATURE_MISMATCH", "Office 파일 형식이 일치하지 않습니다.", 422)


def decode_text(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "cp949"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise AppError("FILE_ENCODING_UNSUPPORTED", "텍스트 인코딩을 확인할 수 없습니다.", 422)


class PlainTextParser:
    def parse(self, filename: str, data: bytes) -> ParsedDocument:
        text = decode_text(data).strip()
        blocks = tuple(
            ExtractedBlock(text=part.strip(), locator={"paragraph": index + 1})
            for index, part in enumerate(re.split(r"\n\s*\n", text))
            if part.strip()
        )
        return ParsedDocument(PurePath(filename).stem, blocks, "plain-text")


class CsvParser:
    def parse(self, filename: str, data: bytes) -> ParsedDocument:
        rows = csv.reader(io.StringIO(decode_text(data)))
        blocks = tuple(
            ExtractedBlock(
                text=" | ".join(cell.strip() for cell in row),
                locator={"row": row_number},
                block_type="table_row",
            )
            for row_number, row in enumerate(rows, start=1)
            if any(cell.strip() for cell in row)
        )
        return ParsedDocument(PurePath(filename).stem, blocks, "csv")


class PdfParser:
    def parse(self, filename: str, data: bytes) -> ParsedDocument:
        reader = PdfReader(io.BytesIO(data))
        blocks = tuple(
            ExtractedBlock(text=text, locator={"page": page_number})
            for page_number, page in enumerate(reader.pages, start=1)
            if (text := (page.extract_text() or "").strip())
        )
        quality = len(blocks) / max(len(reader.pages), 1)
        return ParsedDocument(PurePath(filename).stem, blocks, "pypdf", quality_score=quality)


class DocxParser:
    def parse(self, filename: str, data: bytes) -> ParsedDocument:
        document = DocxDocument(io.BytesIO(data))
        blocks = tuple(
            ExtractedBlock(text=text, locator={"paragraph": index})
            for index, paragraph in enumerate(document.paragraphs, start=1)
            if (text := paragraph.text.strip())
        )
        return ParsedDocument(PurePath(filename).stem, blocks, "python-docx")


class PptxParser:
    def parse(self, filename: str, data: bytes) -> ParsedDocument:
        presentation = Presentation(io.BytesIO(data))
        extracted: list[ExtractedBlock] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            for shape_number, shape in enumerate(slide.shapes, start=1):
                text = getattr(shape, "text", "").strip()
                if text:
                    extracted.append(
                        ExtractedBlock(
                            text=text,
                            locator={"slide": slide_number, "shape": shape_number},
                        )
                    )
        return ParsedDocument(PurePath(filename).stem, tuple(extracted), "python-pptx")


PARSERS: dict[str, Parser] = {
    ".txt": PlainTextParser(),
    ".md": PlainTextParser(),
    ".csv": CsvParser(),
    ".pdf": PdfParser(),
    ".docx": DocxParser(),
    ".pptx": PptxParser(),
}


class _ReadableHtmlParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.blocks: list[ExtractedBlock] = []
        self.title = ""
        self._skip_depth = 0
        self._current_tag = ""
        self._text: list[str] = []

    def handle_starttag(self, tag: str, _attrs: list[tuple[str, str | None]]) -> None:
        tag = tag.casefold()
        if tag in {"script", "style", "noscript", "template", "svg"}:
            self._skip_depth += 1
        if not self._skip_depth and tag in {
            "title",
            "h1",
            "h2",
            "h3",
            "h4",
            "h5",
            "h6",
            "p",
            "li",
            "blockquote",
            "figcaption",
            "td",
            "th",
        }:
            self._flush()
            self._current_tag = tag

    def handle_endtag(self, tag: str) -> None:
        tag = tag.casefold()
        if tag in {"script", "style", "noscript", "template", "svg"}:
            self._skip_depth = max(0, self._skip_depth - 1)
        elif not self._skip_depth and tag == self._current_tag:
            self._flush()

    def handle_data(self, data: str) -> None:
        if not self._skip_depth and self._current_tag:
            self._text.append(data)

    def close(self) -> None:
        super().close()
        self._flush()

    def _flush(self) -> None:
        text = re.sub(r"\s+", " ", " ".join(self._text)).strip()
        if text:
            if self._current_tag == "title" and not self.title:
                self.title = text
            else:
                self.blocks.append(
                    ExtractedBlock(
                        text=text,
                        locator={"element": self._current_tag or "text", "index": len(self.blocks) + 1},
                        block_type=("heading" if self._current_tag.startswith("h") else self._current_tag or "text"),
                    )
                )
        self._text.clear()
        self._current_tag = ""


def _json_blocks(value: object, path: str = "$") -> list[ExtractedBlock]:
    blocks: list[ExtractedBlock] = []
    if isinstance(value, dict):
        for key, item in value.items():
            blocks.extend(_json_blocks(item, f"{path}.{key}"))
    elif isinstance(value, list):
        for index, item in enumerate(value):
            blocks.extend(_json_blocks(item, f"{path}[{index}]"))
    elif value is not None:
        blocks.append(
            ExtractedBlock(
                text=str(value),
                locator={"json_path": path},
                block_type="field",
            )
        )
    return blocks


def parse_fetched_document(url: str, content_type: str, data: bytes, max_size: int) -> ParsedDocument:
    """Parse an approved network response while retaining stable source locators."""
    path = unquote(urlsplit(url).path)
    filename = PurePath(path).name or "source"
    suffix = PurePath(filename).suffix.lower()
    if suffix not in PARSERS:
        suffix = next(
            (candidate for candidate, types in ALLOWED_UPLOADS.items() if content_type in types),
            "",
        )
        if suffix:
            filename = f"{filename}{suffix}"
    if suffix in PARSERS and content_type in ALLOWED_UPLOADS[suffix]:
        return parse_document(filename, content_type, data, max_size)
    if content_type in {"text/html", "application/xhtml+xml"}:
        parser = _ReadableHtmlParser()
        parser.feed(decode_text(data))
        parser.close()
        if not parser.blocks:
            raise AppError("DOCUMENT_TEXT_EMPTY", "웹 문서에서 본문을 추출할 수 없습니다.", 422)
        return ParsedDocument(parser.title or filename, tuple(parser.blocks), "html-parser")
    if content_type in {"application/json", "application/ld+json"}:
        try:
            value = json.loads(decode_text(data))
        except json.JSONDecodeError as exc:
            raise AppError("DOCUMENT_PARSE_FAILED", "JSON 소스를 파싱할 수 없습니다.", 422) from exc
        blocks = tuple(_json_blocks(value))
        if not blocks:
            raise AppError("DOCUMENT_TEXT_EMPTY", "JSON 소스에 처리할 값이 없습니다.", 422)
        return ParsedDocument(filename, blocks, "json")
    if content_type in {"application/xml", "text/xml", "application/rss+xml", "application/atom+xml"}:
        try:
            root = ElementTree.fromstring(data)
        except ElementTree.ParseError as exc:
            raise AppError("DOCUMENT_PARSE_FAILED", "XML 소스를 파싱할 수 없습니다.", 422) from exc
        blocks = tuple(
            ExtractedBlock(
                text=text,
                locator={"element": element.tag.rsplit("}", 1)[-1], "index": index},
                block_type="xml_element",
            )
            for index, element in enumerate(root.iter(), start=1)
            if (text := (element.text or "").strip())
        )
        if not blocks:
            raise AppError("DOCUMENT_TEXT_EMPTY", "XML 소스에 처리할 본문이 없습니다.", 422)
        return ParsedDocument(filename, blocks, "xml")
    if content_type.startswith("text/"):
        return PlainTextParser().parse(filename, data)
    raise AppError("SOURCE_CONTENT_TYPE_BLOCKED", "지원하지 않는 원격 문서 형식입니다.", 422)


def parse_document(filename: str, content_type: str, data: bytes, max_size: int) -> ParsedDocument:
    suffix = validate_upload_metadata(filename, content_type, len(data), max_size)
    validate_file_signature(suffix, data)
    parsed = PARSERS[suffix].parse(filename, data)
    if not parsed.blocks:
        raise AppError(
            "DOCUMENT_TEXT_EMPTY",
            "추출 가능한 텍스트가 없습니다. OCR 검토가 필요합니다.",
            422,
        )
    return parsed


def mask_pii(text: str) -> tuple[str, bool]:
    masked = _EMAIL.sub("[EMAIL]", text)
    masked = _PHONE.sub("[PHONE]", masked)
    masked = _RESIDENT_ID.sub("[IDENTIFIER]", masked)
    return masked, masked != text


def semantic_chunks(
    blocks: tuple[ExtractedBlock, ...], *, max_characters: int = 1_500
) -> tuple[ChunkDraft, ...]:
    if max_characters < 200:
        raise ValueError("max_characters must be at least 200")
    chunks: list[ChunkDraft] = []
    buffer: list[str] = []
    locators: list[dict[str, int | str]] = []

    def flush() -> None:
        if not buffer:
            return
        raw_text = "\n\n".join(buffer).strip()
        text, was_masked = mask_pii(raw_text)
        chunks.append(
            ChunkDraft(
                sequence=len(chunks),
                text=text,
                locator={"start": locators[0], "end": locators[-1]},
                text_hash=hashlib.sha256(text.encode("utf-8")).hexdigest(),
                token_estimate=max(1, len(text) // 3),
                pii_masked=was_masked,
            )
        )
        buffer.clear()
        locators.clear()

    for block in blocks:
        pieces = [block.text[index : index + max_characters] for index in range(0, len(block.text), max_characters)]
        for piece in pieces:
            projected = sum(len(item) for item in buffer) + len(piece) + max(0, len(buffer) * 2)
            if buffer and projected > max_characters:
                flush()
            buffer.append(piece)
            locators.append(block.locator)
    flush()
    return tuple(chunks)
